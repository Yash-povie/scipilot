import os
import io
import uuid
from fastapi import FastAPI
from pydantic import BaseModel
from typing import List, Dict, Any
from opentelemetry.instrumentation.fastapi import FastAPIInstrumentor
from minio import Minio
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from contextlib import asynccontextmanager

MINIO_ENDPOINT = os.getenv("MINIO_ENDPOINT", "minio:9000")
MINIO_ACCESS_KEY = os.getenv("MINIO_ACCESS_KEY", "minioadmin")
MINIO_SECRET_KEY = os.getenv("MINIO_SECRET_KEY", "minioadmin")
BUCKET_NAME = "scipilot"

minio_client = Minio(
    MINIO_ENDPOINT,
    access_key=MINIO_ACCESS_KEY,
    secret_key=MINIO_SECRET_KEY,
    secure=False
)

@asynccontextmanager
async def lifespan(app: FastAPI):
    if not minio_client.bucket_exists(BUCKET_NAME):
        minio_client.make_bucket(BUCKET_NAME)
    yield

app = FastAPI(title="report-service", lifespan=lifespan)

class HealthResponse(BaseModel):
    status: str
    service: str

class PdfRequest(BaseModel):
    markdown: str

class PptxRequest(BaseModel):
    findings: List[Dict[str, Any]]

class ReportResponse(BaseModel):
    url: str

@app.get("/health", response_model=HealthResponse)
async def health_check() -> HealthResponse:
    return HealthResponse(status="ok", service="report-service")

@app.post("/generate/pdf", response_model=ReportResponse)
async def generate_pdf(req: PdfRequest) -> ReportResponse:
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    textobject = c.beginText()
    textobject.setTextOrigin(40, 750)
    textobject.setFont("Helvetica", 12)
    
    for line in req.markdown.split("\n"):
        # Very basic wrapping for demonstration
        textobject.textLine(line[:100])
        if len(line) > 100:
            textobject.textLine(line[100:200])
            
    c.drawText(textobject)
    c.save()
    buffer.seek(0)
    
    filename = f"report_{uuid.uuid4().hex}.pdf"
    minio_client.put_object(
        BUCKET_NAME, 
        filename, 
        buffer, 
        length=buffer.getbuffer().nbytes,
        content_type="application/pdf"
    )
    
    url = minio_client.presigned_get_object(BUCKET_NAME, filename)
    return ReportResponse(url=url)

@app.post("/generate/pptx", response_model=ReportResponse)
async def generate_pptx(req: PptxRequest) -> ReportResponse:
    prs = Presentation()
    
    for finding in req.findings:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = finding.get("title", "Finding")
        content.text = finding.get("description", "")
        
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    filename = f"presentation_{uuid.uuid4().hex}.pptx"
    minio_client.put_object(
        BUCKET_NAME, 
        filename, 
        buffer, 
        length=buffer.getbuffer().nbytes,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    
    url = minio_client.presigned_get_object(BUCKET_NAME, filename)
    return ReportResponse(url=url)

FastAPIInstrumentor.instrument_app(app)
